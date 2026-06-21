"""
build_slides.py
Gera a apresentacao do Tema 4 - Process Mining (Grupo G5).
Rodar de dentro de /code: python build_slides.py
Saída: ../Apresentacao_Tema4_ProcessMining_G5.pptx
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── Paleta ──────────────────────────────────────────────────────────────────
BLUE   = RGBColor(0x1e, 0x49, 0x65)   # UNIFESSPA azul escuro
BLUE2  = RGBColor(0x2a, 0x6f, 0x97)   # azul médio
TEAL   = RGBColor(0x52, 0xb6, 0x9a)   # verde-teal
ORANGE = RGBColor(0xf3, 0x72, 0x2c)   # laranja destaque
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xf0, 0xf4, 0xf8)   # fundo claro
DARK   = RGBColor(0x22, 0x22, 0x33)   # texto escuro
GRAY   = RGBColor(0x88, 0x88, 0x99)   # cinza suave

FIG = os.path.join("..", "report", "figures")
OUT = os.path.join("..", "Apresentacao_Tema4_ProcessMining_G5.pptx")

W, H = Cm(33.87), Cm(19.05)   # 16:9

# ── Helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    # remove todos os layouts desnecessários
    return prs

def blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def rgb_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    rgb_fill(shape, color)
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_para(tf, text, size=16, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    from pptx.util import Pt as PPt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def header_bar(slide, title, subtitle="", presenter=""):
    """Barra azul no topo com título."""
    add_rect(slide, 0, 0, 33.87, 3.2, BLUE)
    add_text(slide, title, 0.8, 0.3, 28, 1.5,
             size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.8, 1.75, 22, 1.0,
                 size=14, color=RGBColor(0xb0, 0xd0, 0xe8))
    if presenter:
        add_text(slide, presenter, 24, 0.5, 9, 2.0,
                 size=11, color=TEAL, align=PP_ALIGN.RIGHT, italic=True)
    # linha decorativa
    add_rect(slide, 0, 3.2, 33.87, 0.12, TEAL)

def add_image(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, Cm(x), Cm(y), Cm(w), Cm(h))
    else:
        slide.shapes.add_picture(path, Cm(x), Cm(y), Cm(w))

def bullet_box(slide, items, x, y, w, h, size=15, title=None, title_color=BLUE2):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = title
        r.font.size = Pt(size + 2)
        r.font.bold = True
        r.font.color.rgb = title_color
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = ("  " if not item.startswith("  ") else "") + item
        r.font.size = Pt(size)
        r.font.color.rgb = DARK

def footer(slide, num):
    add_text(slide, f"Process Mining — Tema 4 | Grupo G5 — UNIFESSPA",
             0.5, 18.2, 26, 0.7, size=9, color=GRAY, italic=True)
    add_text(slide, str(num), 32, 18.2, 1.5, 0.7,
             size=9, color=GRAY, align=PP_ALIGN.RIGHT)

# ── Slide factories ──────────────────────────────────────────────────────────

def slide_capa(prs):
    sl = blank(prs)
    # fundo azul escuro
    add_rect(sl, 0, 0, 33.87, 19.05, BLUE)
    # faixa lateral esquerda
    add_rect(sl, 0, 0, 0.6, 19.05, TEAL)
    # faixa inferior
    add_rect(sl, 0, 15.5, 33.87, 3.55, RGBColor(0x12, 0x2d, 0x42))

    add_text(sl, "UNIVERSIDADE FEDERAL DO SUL E SUDESTE DO PARÁ — UNIFESSPA",
             1.2, 0.6, 31, 1.1, size=11, color=TEAL, bold=True)
    add_text(sl, "Disciplina: Mineração de Dados  |  Prof. Adam D. F. dos Santos",
             1.2, 1.5, 31, 0.9, size=11, color=RGBColor(0x9, 0xb8, 0xcc))

    add_text(sl, "TEMA 4", 1.2, 3.2, 31, 1.4,
             size=14, bold=True, color=TEAL, italic=True)
    add_text(sl, "Process Mining", 1.2, 4.3, 31, 2.2,
             size=40, bold=True, color=WHITE)
    add_text(sl, "Descoberta de Processos e Conformance Checking\nAplicados ao BPI Challenge 2017",
             1.2, 6.8, 31, 1.8, size=18, color=RGBColor(0xb0, 0xd0, 0xe8))

    add_text(sl, "Grupo G5", 1.2, 9.5, 20, 1.0, size=14, bold=True, color=WHITE)
    add_text(sl, "Michelangelo  ·  João Marcos  ·  André  ·  Daniel",
             1.2, 10.3, 31, 0.9, size=14, color=RGBColor(0xb0, 0xd0, 0xe8))

    add_text(sl, "Marabá — PA, junho de 2026", 1.2, 16.1, 20, 0.9,
             size=12, color=GRAY)
    add_text(sl, "1", 32, 16.1, 1.5, 0.9, size=12, color=GRAY, align=PP_ALIGN.RIGHT)

def slide_sumario(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Sumário", "Tópicos abordados na apresentação")
    footer(sl, 2)

    topicos = [
        (BLUE,   "1. Introdução ao Process Mining",              "Contexto, conceitos fundamentais e motivação"),
        (BLUE,   "2. Dataset: BPI Challenge 2017",               "Processo de concessão de empréstimos e adaptação metodológica"),
        (BLUE,   "3. Metodologia e Implementação",               "Pipeline de 7 etapas — pm_toolkit.py"),
        (BLUE2,  "4. Análise Exploratória (EDA)",                "Frequência, duração dos casos e variantes"),
        (BLUE2,  "5. Process Discovery",                         "DFG, Alpha Miner e Inductive Miner"),
        (TEAL,   "6. Conformance Checking",                      "Fitness, Precision, Generalization, Simplicity e desvios"),
        (ORANGE, "7. Performance Analysis e Gargalos",           "Tempo de espera e identificação de bottlenecks"),
        (ORANGE, "8. Síntese, Limitações e Conclusão",           "Insights gerenciais e trabalhos futuros"),
    ]
    for i, (cor, titulo, desc) in enumerate(topicos):
        y = 3.5 + i * 1.85
        add_rect(sl, 1.0, y, 0.35, 1.5, cor)
        add_text(sl, titulo, 1.6, y + 0.05, 20, 0.8, size=14, bold=True, color=cor)
        add_text(sl, desc,   1.6, y + 0.75, 30, 0.7, size=12, color=GRAY)

def slide_section(prs, nome, topico, cor, num):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, BLUE)
    add_rect(sl, 0, 0, 0.6, 19.05, cor)
    add_text(sl, nome, 1.5, 7.0, 30, 2.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, topico, 1.5, 9.5, 30, 1.5, size=18, color=cor, align=PP_ALIGN.CENTER, italic=True)
    footer(sl, num)

def slide_intro(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "O que é Process Mining?", "Extraindo conhecimento real de processos a partir de dados")
    footer(sl, 3)

    # analogia central
    add_rect(sl, 1.0, 3.5, 31.5, 2.2, BLUE)
    add_text(sl, "Ideia central: cada vez que alguém executa uma tarefa no sistema, fica registrado — quem fez, o quê e quando.",
             1.3, 3.6, 31.0, 1.0, size=14, bold=True, color=WHITE)
    add_text(sl, "Process Mining analisa esses registros para revelar como o processo realmente acontece, não como foi planejado.",
             1.3, 4.35, 31.0, 1.0, size=13, color=RGBColor(0xb0, 0xd0, 0xe8))

    # caixa esquerda
    add_rect(sl, 1.0, 6.2, 15.2, 6.0, WHITE)
    bullet_box(sl, [
        "Empresas desenham processos ideais no papel",
        "Na prática, os funcionários seguem caminhos diferentes",
        "Ninguém percebe — até aparecer um gargalo ou erro",
        "",
        "Process Mining usa os logs do sistema para mostrar",
        "o processo real, com todos os desvios e atrasos",
    ], 1.3, 6.4, 14.5, 5.6, size=13,
       title="O Problema", title_color=BLUE)

    # caixa direita
    add_rect(sl, 17.5, 6.2, 15.2, 6.0, WHITE)
    bullet_box(sl, [
        "Process Discovery",
        "  Descobrir o modelo real a partir do log",
        "Conformance Checking",
        "  Comparar o modelo real com o planejado",
        "Performance Analysis",
        "  Medir tempos e identificar gargalos",
    ], 17.8, 6.4, 14.5, 5.6, size=13,
       title="3 Subtarefas Centrais", title_color=TEAL)

    # fluxo base
    for i, (label, cor) in enumerate([
        ("Event Log\n(dados brutos)", BLUE),
        ("Algoritmos\nde PM", BLUE2),
        ("Modelo do\nProcesso Real", TEAL),
    ]):
        x = 1.0 + i * 10.8
        add_rect(sl, x, 13.0, 9.8, 2.8, cor)
        add_text(sl, label, x+0.3, 13.3, 9.2, 2.0, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "→", 11.0, 13.8, 1.5, 1.2, size=22, bold=True, color=BLUE2, align=PP_ALIGN.CENTER)
    add_text(sl, "→", 21.8, 13.8, 1.5, 1.2, size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

def slide_dataset(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Dataset: BPI Challenge 2017", "Adaptação metodológica por restrição de rede")
    footer(sl, 4)

    # coluna esquerda — dataset real
    add_rect(sl, 1.0, 3.6, 14.5, 8.0, WHITE)
    bullet_box(sl, [
        "Log real de concessão de empréstimos",
        "Instituição financeira holandesa",
        "31.509 casos | 1.202.267 eventos",
        "26 atividades — 3 sub-processos:",
        "  Application / Offer / Workflow",
        "Dataset padrão acadêmico (4TU ResearchData)",
    ], 1.3, 3.8, 13.8, 7.5, size=13,
       title="BPI Challenge 2017 (original)", title_color=BLUE)

    # coluna direita — log sintético
    add_rect(sl, 17.0, 3.6, 15.5, 8.0, WHITE)
    bullet_box(sl, [
        "Sem acesso à internet no ambiente usado",
        "Impossível instalar PM4Py ou baixar .xes",
        "Solução: log sintético deterministico (seed=42)",
        "3.000 casos | 54.597 eventos | 23 atividades",
        "Preserva: vocabulário, laços, ramificações",
        "Validade estrutural, não quantitativa absoluta",
    ], 17.3, 3.8, 14.8, 7.5, size=13,
       title="Adaptação Adotada", title_color=ORANGE)

    # faixa de alerta
    add_rect(sl, 1.0, 12.2, 31.5, 1.5, RGBColor(0xff, 0xf0, 0xd0))
    add_text(sl, "Todas as limitações desta adaptação são documentadas no relatório (Seção 6 — Limitações).",
             1.3, 12.35, 31.0, 1.1, size=13, color=RGBColor(0x8, 0x50, 0x20), italic=True)

def slide_metodologia(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Metodologia", "Pipeline em 7 etapas")
    footer(sl, 5)

    etapas = [
        ("1", "Geração do log sintético",      BLUE,   1.0,  3.6),
        ("2", "Importação e limpeza",           BLUE2,  5.6,  3.6),
        ("3", "Análise Exploratória (EDA)",     TEAL,   10.2, 3.6),
        ("4", "Process Discovery",              BLUE,   14.8, 3.6),
        ("5", "Conformance Checking",           BLUE2,  19.4, 3.6),
        ("6", "Performance e Gargalos",         TEAL,   24.0, 3.6),
        ("7", "Síntese e Interpretação",        ORANGE, 28.6, 3.6),
    ]
    for num, label, cor, x, y in etapas:
        add_rect(sl, x, y, 4.1, 4.2, cor)
        add_text(sl, num, x+0.2, y+0.15, 3.5, 1.2, size=28, bold=True, color=WHITE)
        add_text(sl, label, x+0.2, y+1.5, 3.6, 2.3, size=12, bold=True, color=WHITE)

    # setas entre etapas (linha)
    add_rect(sl, 5.1, 5.3, 0.5, 0.3, GRAY)

    # ferramentas
    add_rect(sl, 1.0, 9.5, 31.5, 4.5, WHITE)
    bullet_box(sl, [
        "Python 3  |  pandas / numpy  |  matplotlib / seaborn  |  graphviz",
        "pm_toolkit.py — biblioteca própria: DFG, Alpha Miner, Inductive Miner,",
        "  conformance checking e performance analysis implementados do zero",
        "Objetivo: independência de bibliotecas externas (PM4Py indisponível)",
    ], 1.3, 9.6, 31.0, 4.2, size=13, title="Ferramentas e Implementação", title_color=BLUE2)

def slide_eda1(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Análise Exploratória — Visão Geral", "3.000 casos · 54.597 eventos · 23 atividades")
    footer(sl, 6)

    # KPIs
    kpis = [
        ("3.000", "casos",         BLUE),
        ("54.597", "eventos",      BLUE2),
        ("18,20", "eventos/caso",  TEAL),
        ("169", "variantes",       ORANGE),
    ]
    for i, (val, lab, cor) in enumerate(kpis):
        x = 1.0 + i * 8.0
        add_rect(sl, x, 3.6, 7.2, 2.6, cor)
        add_text(sl, val, x+0.3, 3.7, 6.5, 1.5, size=28, bold=True, color=WHITE)
        add_text(sl, lab, x+0.3, 5.0, 6.5, 0.9, size=13, color=WHITE)

    # figuras lado a lado
    add_image(sl, os.path.join(FIG, "01_activity_frequency.png"), 0.8, 6.7, 15.5)
    add_image(sl, os.path.join(FIG, "02_temporal_distribution.png"), 17.2, 6.7, 15.5)

    add_text(sl, "Frequência de atividades", 0.8, 16.3, 15.5, 0.8,
             size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_text(sl, "Distribuição temporal dos eventos", 17.2, 16.3, 15.5, 0.8,
             size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

def slide_eda2(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Análise Exploratória — Duração e Variantes", "Lead time e distribuição de variantes de execução")
    footer(sl, 7)

    add_image(sl, os.path.join(FIG, "03_case_duration_distribution.png"), 0.8, 3.6, 15.5)
    add_image(sl, os.path.join(FIG, "05_top_variants.png"), 17.2, 3.6, 15.5)

    add_rect(sl, 0.8, 14.8, 15.2, 3.4, WHITE)
    bullet_box(sl, [
        "Duração média: 4,82 dias (mediana 4,14 d)",
        "P90: 8,43 d | Máximo: 39,17 d",
        "Cauda longa indica laços de re-trabalho",
    ], 1.0, 14.9, 14.8, 3.1, size=13)

    add_rect(sl, 17.2, 14.8, 15.5, 3.4, WHITE)
    bullet_box(sl, [
        "169 variantes distintas de execução",
        "Top 1 variante: 19,33% dos casos",
        "Top 5 variantes: 44,47% dos casos",
        "Distribuição de cauda longa típica",
    ], 17.4, 14.9, 15.1, 3.1, size=13)

def slide_dfg(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Process Discovery — DFG", "Directly-Follows Graph: cada seta = uma sequência observada no log")
    footer(sl, 8)

    add_image(sl, os.path.join(FIG, "06_dfg_full.png"), 0.5, 3.4, 16.0)
    add_image(sl, os.path.join(FIG, "07_dfg_filtered_heuristic.png"), 17.2, 3.4, 16.0)

    add_rect(sl, 0.5, 15.5, 16.0, 3.0, WHITE)
    add_text(sl, "DFG Completo", 0.7, 15.6, 15.6, 0.8, size=13, color=BLUE, bold=True)
    add_text(sl, "Todas as sequências observadas, incluindo raras.",
             0.7, 16.2, 15.6, 0.7, size=12, color=DARK)
    add_text(sl, "Setas mais grossas = sequências mais frequentes.",
             0.7, 16.85, 15.6, 0.7, size=12, color=DARK, italic=True)

    add_rect(sl, 17.2, 15.5, 15.5, 3.0, WHITE)
    add_text(sl, "DFG Filtrado (≥10% de frequência)", 17.4, 15.6, 15.1, 0.8, size=13, color=TEAL, bold=True)
    add_text(sl, "Só o caminho principal — como a maioria dos casos flui.",
             17.4, 16.2, 15.1, 0.7, size=12, color=DARK)
    add_text(sl, "Equivale ao Heuristic Miner: ignora sequências raras.",
             17.4, 16.85, 15.1, 0.7, size=12, color=DARK, italic=True)

def slide_alpha(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Alpha Miner — Rede de Petri", "Modelo formal com lugares (condições) e transições (atividades)")
    footer(sl, 9)

    add_image(sl, os.path.join(FIG, "08_alpha_miner_net.png"), 0.5, 3.4, 22.0)

    add_rect(sl, 23.2, 3.4, 9.8, 10.5, WHITE)
    bullet_box(sl, [
        "Quadrados laranja = atividades",
        "Círculos verdes = condições",
        "  (ex.: 'oferta enviada')",
        "",
        "O algoritmo analisa:",
        "  Quem vem antes de quem?",
        "  Quem ocorre em paralelo?",
        "  Quem nunca ocorre junto?",
        "",
        "Ponto fraco: não lida bem com",
        "laços curtos e repetições",
    ], 23.4, 3.6, 9.2, 10.0, size=12, title="Como ler o grafo", title_color=BLUE)

    add_rect(sl, 0.5, 14.6, 32.5, 3.8, RGBColor(0xe8, 0xf4, 0xf8))
    add_text(sl, "Resultado para o BPI17:",
             0.8, 14.7, 15.0, 0.7, size=13, bold=True, color=BLUE)
    add_text(sl, "23 atividades gerou 23 condições — estrutura densa, reflexo dos laços de re-trabalho.",
             0.8, 15.3, 31.5, 0.7, size=12, color=DARK)
    add_text(sl, "Fitness: 1,000  |  Precision: 1,000  |  Simplicity: 0,906  — melhor equilíbrio entre os modelos testados.",
             0.8, 16.0, 31.5, 0.8, size=12, color=DARK, bold=True)

def slide_conformance(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Conformance Checking — Resultados", "Comparação entre os modelos descobertos")
    footer(sl, 10)

    # tabela manual
    rows = [
        ("Modelo",                     "Fitness", "Precision", "Generaliz.", "Simplicity", "Arestas", True),
        ("DFG filtrado (Heuristic)",   "0,9893",  "1,0000",    "1,0000",     "0,9149",     "24",      False),
        ("Alpha Miner",                "1,0000",  "1,0000",    "1,0000",     "0,9058",     "29",      False),
        ("Inductive Miner",            "0,9973",  "0,0664",    "0,0664",     "0,1938",     "422",     False),
        ("DFG completo (bruto)",       "1,0000",  "1,0000",    "1,0000",     "0,9058",     "29",      False),
    ]
    col_x = [1.0, 11.0, 16.5, 20.5, 24.5, 28.5]
    col_w = [9.5, 5.0, 4.0, 4.0, 4.0, 4.5]
    for ri, row in enumerate(rows):
        y = 3.6 + ri * 1.6
        bg = BLUE if row[6] else (RGBColor(0xe8, 0xf4, 0xf8) if ri % 2 == 0 else WHITE)
        add_rect(sl, 0.8, y, 31.8, 1.5, bg)
        for ci, (val, cx, cw) in enumerate(zip(row[:6], col_x, col_w)):
            tc = WHITE if row[6] else (ORANGE if ci > 0 and ri == 3 and val not in ("Inductive Miner","0,9973") else DARK)
            add_text(sl, val, cx, y+0.2, cw, 1.1,
                     size=12, bold=row[6], color=tc)

    # destaque visual do Inductive Miner
    add_rect(sl, 0.8, 9.6, 31.8, 1.5, RGBColor(0xff, 0xf0, 0xd8))
    add_text(sl, "⚠  Inductive Miner caiu no 'modelo flor': permite qualquer sequência → fitness alta, mas precision quase zero.",
             1.1, 9.75, 31.2, 1.1, size=12, color=RGBColor(0x8B, 0x45, 0x13), bold=True)

    add_rect(sl, 0.8, 11.4, 31.8, 5.5, WHITE)
    bullet_box(sl, [
        "Fitness mede: 'o modelo consegue reproduzir o que aconteceu no log?'",
        "Precision mede: 'o modelo só permite o que realmente acontece?'",
        "",
        "Inductive Miner (fitness=0,997 / precision=0,066):",
        "  Aceita quase tudo — inclusive sequências impossíveis na prática",
        "Alpha Miner (fitness=1,0 / precision=1,0 / 29 arestas):",
        "  Modelo mais útil — representa exatamente o que foi observado",
    ], 1.0, 11.5, 31.4, 5.2, size=12)

def slide_desvios(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Conformance Checking — Desvios", "Top desvios contra o modelo DFG filtrado")
    footer(sl, 11)

    add_image(sl, os.path.join(FIG, "09_conformance_metrics.png"), 0.5, 3.5, 18.0)
    add_image(sl, os.path.join(FIG, "10_top_deviations.png"), 19.0, 3.5, 14.0)

    add_rect(sl, 0.5, 15.2, 32.5, 3.0, WHITE)
    bullet_box(sl, [
        "Laço de múltiplas ofertas (267 ocorr.): estruturalmente esperado no processo de empréstimos",
        "Omissão de W_Handle leads (171): via alternativa conhecida",
        "Ramificação de análise de fraude (126/78/48): exceção rara, mas processualmente distinta",
    ], 0.7, 15.3, 32.0, 2.7, size=12)

def slide_performance(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Performance Analysis — Gargalos", "Tempo de espera entre atividades consecutivas")
    footer(sl, 12)

    add_image(sl, os.path.join(FIG, "11_bottlenecks.png"), 0.5, 3.5, 18.5)
    add_image(sl, os.path.join(FIG, "12_wait_by_activity.png"), 19.5, 3.5, 13.5)

    add_rect(sl, 0.5, 15.2, 31.5, 3.5, WHITE)
    add_text(sl, "O que os dados revelam:",
             0.8, 15.3, 20.0, 0.7, size=13, bold=True, color=BLUE)
    bullet_box(sl, [
        "Maior gargalo: após enviar a oferta, o cliente demora em média 1,27 dias para ser contatado",
        "Todos os atrasos estão em etapas humanas — ligações, validação de documentos",
        "Etapas automáticas (criar oferta, enviar email) levam segundos",
        "Conclusão prática: mais atendentes ou automação da triagem reduziria o tempo total do processo",
    ], 0.8, 15.9, 31.2, 2.7, size=12)

def slide_sintese(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Análise Crítica — O que os resultados significam?", "Além dos números: interpretação e implicações")
    footer(sl, 13)

    insights = [
        (BLUE,
         "O processo real é muito mais caótico que o planejado",
         "169 variantes — só 19,3% seguem o fluxo ideal. Os outros 80,7% desviam por re-trabalho, múltiplas ofertas ou fraude.",
         "Por que isso importa? Um gestor que só olha o fluxograma não enxerga a complexidade real do dia a dia."),
        (TEAL,
         "Os gargalos são humanos, não tecnológicos",
         "Todas as etapas automáticas (criar oferta, enviar email) levam segundos. Os atrasos estão nas ligações e validações.",
         "Por que isso importa? Investir em mais servidores não resolve. O problema é de capacidade humana e triagem."),
        (ORANGE,
         "Alta fitness não garante modelo útil — cuidado com o modelo flor",
         "Inductive Miner: fitness 99,7% mas precision 6,6%. O modelo aceita praticamente qualquer sequência de atividades.",
         "Por que isso importa? Uma métrica boa isolada pode enganar. É preciso analisar o conjunto fitness + precision."),
        (BLUE2,
         "A escolha do algoritmo depende do objetivo",
         "Alpha Miner revelou o processo de forma compacta (29 arestas). Inductive Miner é mais formal, mas menos útil aqui.",
         "Por que isso importa? Não existe algoritmo universalmente melhor — o contexto define a escolha."),
    ]
    for i, (cor, titulo, texto, implicacao) in enumerate(insights):
        y = 3.5 + i * 3.7
        add_rect(sl, 0.8, y, 0.4, 3.2, cor)
        add_text(sl, titulo, 1.4, y+0.05, 31.5, 0.9, size=13, bold=True, color=cor)
        add_text(sl, texto, 1.4, y+0.9, 31.5, 0.9, size=12, color=DARK)
        add_text(sl, implicacao, 1.4, y+1.75, 31.5, 0.9, size=11, color=GRAY, italic=True)

def slide_limitacoes(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, LIGHT)
    header_bar(sl, "Limitações — Honestidade Metodológica", "O que afeta os resultados e como mitigar")
    footer(sl, 14)

    lims = [
        (ORANGE, "Log sintético, não o BPI17 real",
                 "Sem internet no ambiente de execução — impossível baixar o arquivo original (~1 GB).",
                 "Impacto: números absolutos não comparáveis à literatura. Padrões qualitativos preservados."),
        (ORANGE, "Volume reduzido: 3.000 casos vs. 31.509 reais",
                 "Restrição de desempenho computacional do ambiente sem bibliotecas otimizadas.",
                 "Impacto: algumas métricas estatísticas podem ser menos robustas em dados reais."),
        (BLUE2,  "Algoritmos implementados do zero, sem PM4Py",
                 "PM4Py indisponível. Implementações didáticas — não validadas contra benchmark.",
                 "Impacto: valores de fitness/precision são aproximações, não resultados de produção."),
        (TEAL,   "Heuristic Miner não implementado separadamente",
                 "Representado pelo DFG filtrado por frequência — conceito equivalente, não idêntico.",
                 "Trabalho futuro: implementar com métricas de dependência estatística formais."),
    ]
    for i, (cor, titulo, desc, impacto) in enumerate(lims):
        y = 3.5 + i * 3.7
        add_rect(sl, 0.8, y, 0.4, 3.2, cor)
        add_text(sl, titulo, 1.4, y+0.05, 31.5, 0.85, size=13, bold=True, color=cor)
        add_text(sl, desc, 1.4, y+0.85, 31.5, 0.8, size=12, color=DARK)
        add_text(sl, impacto, 1.4, y+1.65, 31.5, 0.9, size=11, color=GRAY, italic=True)

def slide_conclusao(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, BLUE)
    add_rect(sl, 0, 0, 0.6, 19.05, TEAL)
    add_rect(sl, 0, 15.5, 33.87, 3.55, RGBColor(0x12, 0x2d, 0x42))

    add_text(sl, "Conclusão", 1.5, 1.5, 30, 1.5,
             size=30, bold=True, color=WHITE)
    add_rect(sl, 1.2, 3.0, 25, 0.1, TEAL)

    pontos = [
        "Pipeline completo de Process Mining: EDA, descoberta, conformance e performance",
        "169 variantes de execução — processo real substancialmente mais complexo que o ideal",
        "Gargalos em etapas humanas com espera média superior a 25 horas",
        "Trade-off fitness/precision: Inductive Miner (0,997/0,066) vs Alpha Miner equilibrado",
        "Metodologia robusta mesmo sem PM4Py e sem dataset real — limitações explicitadas",
    ]
    for i, p in enumerate(pontos):
        add_rect(sl, 1.2, 3.8 + i * 2.0, 0.4, 1.3, TEAL)
        add_text(sl, p, 2.0, 3.8 + i * 2.0, 30.5, 1.5, size=14, color=WHITE)

    add_text(sl, "Grupo G5  |  Michelangelo · João Marcos · André · Daniel",
             1.5, 16.0, 28, 0.9, size=12, color=GRAY)
    add_text(sl, "15", 32, 16.0, 1.5, 0.9, size=12, color=GRAY, align=PP_ALIGN.RIGHT)

def slide_perguntas(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, 33.87, 19.05, BLUE)
    add_rect(sl, 0, 0, 0.6, 19.05, ORANGE)

    add_text(sl, "Obrigado!", 1.5, 5.0, 30, 2.5,
             size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Perguntas?", 1.5, 8.0, 30, 2.0,
             size=28, color=TEAL, align=PP_ALIGN.CENTER, italic=True)

    add_text(sl, "Grupo G5 — Process Mining — Tema 4 — Mineração de Dados — UNIFESSPA",
             1.5, 13.5, 30, 1.0, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(sl, "Repositório: github.com/Michelangelo-Costa/ProcessMining-Tema4-G5",
             1.5, 14.5, 30, 1.0, size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_text(sl, "16", 32, 17.8, 1.5, 0.9, size=12, color=GRAY, align=PP_ALIGN.RIGHT)

# ── Main ─────────────────────────────────────────────────────────────────────
prs = new_prs()

slide_capa(prs)         # 1
slide_sumario(prs)      # 2
slide_intro(prs)        # 3
slide_dataset(prs)      # 4
slide_metodologia(prs)  # 5
slide_eda1(prs)         # 6
slide_eda2(prs)         # 7
slide_dfg(prs)          # 8
slide_alpha(prs)        # 9
slide_conformance(prs)  # 10
slide_desvios(prs)      # 11
slide_performance(prs)  # 12
slide_sintese(prs)      # 13
slide_limitacoes(prs)   # 14
slide_conclusao(prs)    # 15
slide_perguntas(prs)    # 16

prs.save(OUT)
print(f"Salvo: {OUT}  ({len(prs.slides)} slides)")
